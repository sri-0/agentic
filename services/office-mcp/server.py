"""Office-document MCP server (Phase 06).

A thin Python MCP server the Go gateway connects to as a remote client (streamable
HTTP). It generates Word/PowerPoint/Excel files and returns a URL — the backend
does not share this process's filesystem, so tools upload (here: write to a served
output dir) and return a link rather than a local path.

Run:
    pip install -r requirements.txt
    OFFICE_OUTPUT_DIR=./out OFFICE_BASE_URL=http://localhost:8090/files \\
        python server.py            # serves MCP on :8090/mcp (streamable HTTP)

Register in the gateway's config/<env>/mcp.yaml:
    servers:
      office:
        type: remote
        url: http://localhost:8090/mcp
        enabled: true
"""
from __future__ import annotations

import os
import uuid
from typing import Any

from mcp.server.fastmcp import FastMCP

from starlette.requests import Request
from starlette.responses import FileResponse, PlainTextResponse, Response

from docxtpl import DocxTemplate
from pptx import Presentation
from pptx.util import Inches, Pt
from openpyxl import Workbook

OUTPUT_DIR = os.path.abspath(os.environ.get("OFFICE_OUTPUT_DIR", "./out"))
BASE_URL = os.environ.get("OFFICE_BASE_URL", "http://localhost:8090/files")
TEMPLATE_DIR = os.environ.get("OFFICE_TEMPLATE_DIR", "./templates")
PORT = int(os.environ.get("OFFICE_PORT", "8090"))

os.makedirs(OUTPUT_DIR, exist_ok=True)

mcp = FastMCP("office", host="0.0.0.0", port=PORT)


# Serve generated files at /files/<name> on the SAME ASGI app that mcp.run
# serves (streamable-http). Without this, every URL returned by the tools 404s
# because nothing served OFFICE_OUTPUT_DIR (finding N1). custom_route mounts a
# public route on the FastMCP Starlette app, so both /mcp and /files/<name>
# work on the one port.
@mcp.custom_route("/files/{filename}", methods=["GET"])
async def serve_file(request: Request) -> Response:
    filename = request.path_params["filename"]
    # Prevent path traversal: only serve plain names inside OUTPUT_DIR.
    if "/" in filename or "\\" in filename or filename in ("", ".", ".."):
        return PlainTextResponse("bad filename", status_code=400)
    path = os.path.abspath(os.path.join(OUTPUT_DIR, filename))
    if os.path.commonpath([OUTPUT_DIR, path]) != OUTPUT_DIR or not os.path.isfile(path):
        return PlainTextResponse("not found", status_code=404)
    return FileResponse(path, filename=filename)


def _save_url(filename: str) -> str:
    """Return the public URL for a generated file under OUTPUT_DIR."""
    return f"{BASE_URL.rstrip('/')}/{filename}"


DEFAULT_TEMPLATE = "report.docx"


def _ensure_default_template() -> str:
    """Ensure a minimal default docxtpl template exists; return its path.

    Ships a programmatic fallback so render_report_docx works out of the box
    (the plan's requirement to make docxtpl actually testable). Designers can
    drop a branded report.docx into OFFICE_TEMPLATE_DIR to override it.
    """
    os.makedirs(TEMPLATE_DIR, exist_ok=True)
    path = os.path.join(TEMPLATE_DIR, DEFAULT_TEMPLATE)
    if not os.path.isfile(path):
        from make_template import build_default_template  # local module

        build_default_template(path)
    return path


@mcp.tool()
def render_report_docx(template: str, context: dict[str, Any]) -> str:
    """Render a Word (.docx) report from a docxtpl template + a JSON context.

    `template` is a filename under OFFICE_TEMPLATE_DIR (a real .docx with Jinja2
    placeholders). Pass an empty string to use the built-in default template
    (`report.docx`, with {{ title }}, {{ author }}, {{ summary }} and a
    {% for s in sections %} loop). `context` fills the placeholders. Returns a
    URL to the generated document. Highest fidelity for branded reports.
    """
    if not template:
        tpl_path = _ensure_default_template()
    else:
        tpl_path = os.path.join(TEMPLATE_DIR, template)
        if not os.path.isfile(tpl_path):
            raise FileNotFoundError(f"template not found: {template}")
    doc = DocxTemplate(tpl_path)
    doc.render(context)
    name = f"report-{uuid.uuid4().hex[:10]}.docx"
    doc.save(os.path.join(OUTPUT_DIR, name))
    return _save_url(name)


@mcp.tool()
def create_pptx(title: str, slides: list[dict[str, Any]]) -> str:
    """Create a PowerPoint (.pptx). Each slide is {"heading": str, "bullets": [str]}.

    Returns a URL to the generated deck. (For branded decks, point this at a
    corporate .potx template — see the GongRzhe PowerPoint MCP server.)
    """
    prs = Presentation()
    title_layout = prs.slide_layouts[0]
    s = prs.slides.add_slide(title_layout)
    s.shapes.title.text = title

    bullet_layout = prs.slide_layouts[1]
    for slide in slides:
        sl = prs.slides.add_slide(bullet_layout)
        sl.shapes.title.text = slide.get("heading", "")
        body = sl.placeholders[1].text_frame
        body.clear()
        for i, bullet in enumerate(slide.get("bullets", [])):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.text = str(bullet)
            p.font.size = Pt(18)

    name = f"deck-{uuid.uuid4().hex[:10]}.pptx"
    prs.save(os.path.join(OUTPUT_DIR, name))
    return _save_url(name)


@mcp.tool()
def create_xlsx(sheet_name: str, header: list[str], rows: list[list[Any]]) -> str:
    """Create an Excel (.xlsx) with a header row and data rows. Returns a URL."""
    wb = Workbook()
    ws = wb.active
    ws.title = sheet_name or "Sheet1"
    if header:
        ws.append(list(header))
    for row in rows:
        ws.append(list(row))
    name = f"sheet-{uuid.uuid4().hex[:10]}.xlsx"
    wb.save(os.path.join(OUTPUT_DIR, name))
    return _save_url(name)


if __name__ == "__main__":
    # Streamable HTTP transport so the Go backend connects as an ordinary remote
    # MCP client at http://<host>:<port>/mcp.
    mcp.run(transport="streamable-http")
